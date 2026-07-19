"""Document loaders: file → Document(text + metadata).

All loaders return the same Document shape so the chunking/embedding
pipeline stays format-agnostic. Supported: PDF, DOCX, PPTX, XLSX.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


@dataclass
class Document:
    """A loaded document, normalized across formats."""

    text: str
    metadata: dict = field(default_factory=dict)


class UnsupportedFormatError(ValueError):
    """Raised when a file extension has no registered loader."""


def load_pdf(path: str | Path) -> Document:
    """Extract text from a PDF, page by page.

    Page breaks are preserved as double newlines; per-page text is also
    stored in metadata so citations can reference page numbers later.
    """
    path = Path(path)
    reader = PdfReader(path)

    pages: list[str] = []
    for page in reader.pages:
        pages.append((page.extract_text() or "").strip())

    text = "\n\n".join(p for p in pages if p)
    return Document(
        text=text,
        metadata={
            "source": path.name,
            "filetype": "pdf",
            "num_pages": len(reader.pages),
            "pages": pages,
        },
    )


def load_docx(path: str | Path) -> Document:
    """Extract text from a Word document: paragraphs + table cells."""
    path = Path(path)
    doc = DocxDocument(str(path))

    parts: list[str] = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return Document(
        text="\n\n".join(parts),
        metadata={
            "source": path.name,
            "filetype": "docx",
            "num_paragraphs": len(doc.paragraphs),
            "num_tables": len(doc.tables),
        },
    )


def load_pptx(path: str | Path) -> Document:
    """Extract text from a PowerPoint deck, slide by slide."""
    path = Path(path)
    prs = Presentation(str(path))

    slides: list[str] = []
    for slide in prs.slides:
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
        ]
        slides.append("\n".join(texts))

    return Document(
        text="\n\n".join(s for s in slides if s),
        metadata={
            "source": path.name,
            "filetype": "pptx",
            "num_slides": len(prs.slides),
            "slides": slides,
        },
    )


def load_xlsx(path: str | Path) -> Document:
    """Extract cell values from an Excel workbook, sheet by sheet.

    Rows are rendered as pipe-separated lines so tabular structure
    survives into plain text for chunking/search.
    """
    path = Path(path)
    wb = load_workbook(str(path), read_only=True, data_only=True)

    sheets: list[str] = []
    sheet_names: list[str] = []
    for ws in wb.worksheets:
        sheet_names.append(ws.title)
        rows: list[str] = [f"[Sheet: {ws.title}]"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        sheets.append("\n".join(rows))
    wb.close()

    return Document(
        text="\n\n".join(sheets),
        metadata={
            "source": path.name,
            "filetype": "xlsx",
            "num_sheets": len(sheet_names),
            "sheet_names": sheet_names,
        },
    )


_LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".xlsx": load_xlsx,
}


def load(path: str | Path) -> Document:
    """Load any supported file by extension."""
    path = Path(path)
    loader = _LOADERS.get(path.suffix.lower())
    if loader is None:
        supported = ", ".join(sorted(_LOADERS))
        raise UnsupportedFormatError(
            f"No loader for '{path.suffix}'. Supported: {supported}"
        )
    return loader(path)
