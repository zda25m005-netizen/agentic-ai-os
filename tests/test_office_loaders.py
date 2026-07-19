from pathlib import Path

import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from app.rag.loaders import load, load_docx, load_pptx, load_xlsx


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    doc = DocxDocument()
    doc.add_paragraph("Quarterly revenue grew 12% year over year.")
    doc.add_paragraph("Risks include supply chain delays.")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Region"
    table.rows[0].cells[1].text = "Revenue"
    p = tmp_path / "report.docx"
    doc.save(p)
    return p


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Q3 Strategy Overview"
    p = tmp_path / "deck.pptx"
    prs.save(p)
    return p


@pytest.fixture
def sample_xlsx(tmp_path: Path) -> Path:
    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["Region", "Revenue"])
    ws.append(["APAC", 42000])
    p = tmp_path / "sales.xlsx"
    wb.save(p)
    return p


def test_load_docx(sample_docx: Path):
    doc = load_docx(sample_docx)
    assert "Quarterly revenue" in doc.text
    assert "Region | Revenue" in doc.text
    assert doc.metadata["filetype"] == "docx"
    assert doc.metadata["num_tables"] == 1


def test_load_pptx(sample_pptx: Path):
    doc = load_pptx(sample_pptx)
    assert "Q3 Strategy Overview" in doc.text
    assert doc.metadata["num_slides"] == 1


def test_load_xlsx(sample_xlsx: Path):
    doc = load_xlsx(sample_xlsx)
    assert "[Sheet: Sales]" in doc.text
    assert "APAC | 42000" in doc.text
    assert doc.metadata["sheet_names"] == ["Sales"]


def test_registry_dispatches_all(sample_docx, sample_pptx, sample_xlsx):
    assert load(sample_docx).metadata["filetype"] == "docx"
    assert load(sample_pptx).metadata["filetype"] == "pptx"
    assert load(sample_xlsx).metadata["filetype"] == "xlsx"
